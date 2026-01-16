import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import time
import ai_engine 

st.set_page_config(page_title="智讲 SmartPresenter", layout="wide", page_icon="🧠")

with st.sidebar:
    st.title("🧠 智讲 Pro")
    st.caption("真实模型加载版")
    api_key = st.text_input("🔑 Google API Key", type="password")
    
    available_models = []
    if api_key:
        # 调用引擎获取真实列表
        success, result = ai_engine.configure_genai(api_key)
        if success:
            available_models = result
            st.success(f"✅ 成功加载 {len(available_models)} 个模型")
        else:
            st.error(f"❌ 连接失败: {result}")
    
    # 让用户从真实列表中选择，这样绝对不会 404
    if available_models:
        # 自动帮用户选一个带 flash 的
        default_index = 0
        for i, name in enumerate(available_models):
            if "flash" in name and "1.5" in name:
                default_index = i
                break
        selected_model = st.selectbox("👇 请选择模型:", available_models, index=default_index)
    else:
        selected_model = st.selectbox("模型:", ["models/gemini-1.5-flash"])

# 回调函数
def update_status_ui(slide_index, wait_seconds, attempt, max_retries):
    with st.empty():
        for t in range(wait_seconds, 0, -1):
            st.warning(f"⏳ 第 {slide_index} 页触发限流，冷却中... {t}s (重试 {attempt}/{max_retries})")
            time.sleep(1)

st.title("🎙️ 智讲 SmartPresenter")
uploaded_file = st.file_uploader("上传 PPTX", type=['pptx'])

if 'results' not in st.session_state:
    st.session_state['results'] = []

if uploaded_file and api_key and available_models:
    if st.button("🚀 启动分析"):
        st.session_state['results'] = [] 
        genai.configure(api_key=api_key)
        # 这里的 selected_model 是从 Google 列表里直接拿的，绝对保真
        model = genai.GenerativeModel(selected_model)
        prs = Presentation(uploaded_file)
        
        progress_bar = st.progress(0)
        status_box = st.empty()
        result_area = st.container()
        total = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            idx = i + 1
            status_box.info(f"🚀 正在分析第 {idx}/{total} 页 (使用: {selected_model})...")
            progress_bar.progress(i / total)
            
            try:
                data = ai_engine.analyze_slide_content(model, slide, idx, status_callback=update_status_ui)
                data['index'] = idx
                st.session_state['results'].append(data)
                
                with result_area:
                    with st.expander(f"✅ 第 {idx} 页 | {data.get('visual_summary')}", expanded=True):
                        st.write(data['scripts']['standard'])
                
                # 成功后休息 2 秒
                time.sleep(2)

            except Exception as e:
                st.error(f"第 {idx} 页失败: {e}")
        
        status_box.success("🎉 完成！")
        progress_bar.progress(1.0)

elif st.session_state['results']:
    st.divider()
    for data in st.session_state['results']:
        with st.expander(f"第 {data['index']} 页"):
            st.write(data['scripts']['standard'])
