import streamlit as st
import google.generativeai as genai
from pptx import Presentation
import time
import ai_engine  # 引用刚才那个文件

st.set_page_config(page_title="智讲 SmartPresenter", layout="wide", page_icon="🧠")

# --- 侧边栏 ---
with st.sidebar:
    st.title("🧠 智讲 Pro")
    st.caption("双核架构版")
    api_key = st.text_input("🔑 Google API Key", type="password")
    
    available_models = []
    if api_key:
        # 调用 ai_engine 里的函数
        success, result = ai_engine.configure_genai(api_key)
        if success:
            available_models = result
            st.success(f"✅ 连接成功")
        else:
            st.error(f"❌ 连接失败: {result}")
    
    if available_models:
        default_index = 0
        for i, name in enumerate(available_models):
            if "flash" in name and "1.5" in name:
                default_index = i
                break
        selected_model = st.selectbox("👇 选择模型:", available_models, index=default_index)
    else:
        selected_model = st.selectbox("模型:", ["models/gemini-1.5-flash"])

# --- 回调函数 ---
def update_status_ui(slide_index, wait_seconds, attempt, max_retries):
    with st.empty():
        for t in range(wait_seconds, 0, -1):
            st.warning(f"⏳ 第 {slide_index} 页触发限流，冷却中... {t}s (重试 {attempt}/{max_retries})")
            time.sleep(1)

# --- 主界面 ---
st.title("🎙️ 智讲 SmartPresenter")
uploaded_file = st.file_uploader("上传 PPTX", type=['pptx'])

if 'results' not in st.session_state:
    st.session_state['results'] = []

if uploaded_file and api_key and available_models:
    if st.button("🚀 启动分析"):
        st.session_state['results'] = [] 
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel(selected_model)
        prs = Presentation(uploaded_file)
        
        progress_bar = st.progress(0)
        status_box = st.empty()
        result_area = st.container()
        total = len(prs.slides)
        
        for i, slide in enumerate(prs.slides):
            idx = i + 1
            status_box.info(f"🚀 正在分析第 {idx}/{total} 页...")
            progress_bar.progress(i / total)
            
            try:
                # 调用 ai_engine
                data = ai_engine.analyze_slide_content(model, slide, idx, status_callback=update_status_ui)
                data['index'] = idx
                st.session_state['results'].append(data)
                
                with result_area:
                    with st.expander(f"✅ 第 {idx} 页 | {data.get('visual_summary')}", expanded=True):
                        st.write(data['scripts']['standard'])
                
                time.sleep(2)

            except Exception as e:
                st.error(f"第 {idx} 页失败: {e}")
        
        status_box.success("🎉 完成！")
        progress_bar.progress(1.0)

# ... (前面的代码保持不变)

# ... (前面的代码保持不变，从 elif st.session_state['results']: 开始替换)

elif st.session_state['results']:
    st.divider()
    
    # --- 布局：左侧是功能区，右侧是下载区 ---
    col_action, col_download = st.columns([1, 1])
    
    with col_action:
        # 新功能：召唤全局架构师
        if st.button("🧠 生成全篇逻辑诊断 (Global Review)"):
            with st.spinner("AI 正在通读全篇，寻找逻辑漏洞..."):
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel(selected_model)
                
                # 调用我们在 ai_engine 里新写的函数
                review_data = ai_engine.analyze_presentation_logic(model, st.session_state['results'])
                st.session_state['review'] = review_data
    
    with col_download:
        # 实用功能：一键导出 Markdown
        def generate_markdown(results, review=None):
            md = "# 🎙️ 智讲 SmartPresenter 分析报告\n\n"
            
            if review:
                md += "## 🧠 全局逻辑诊断\n"
                md += f"**摘要**: {review['executive_summary']}\n\n"
                md += f"**逻辑分析**: {review['logic_diagnosis']}\n"
                md += f"**优势**: {review['strengths']}\n"
                md += f"**改进建议**: {review['weaknesses']}\n"
                md += f"**金句结语**: {review['closing_remark']}\n\n"
                md += "---\n\n"
            
            for slide in results:
                md += f"## 第 {slide['index']} 页\n"
                md += f"**画面**: {slide.get('visual_summary')}\n\n"
                md += f"**演讲稿**: {slide['scripts']['standard']}\n\n"
                md += f"> 知识点: {slide['knowledge_extension']['trivia']}\n\n"
                md += "---\n"
            return md

        # 如果生成过报告，把报告也加进下载里；如果没有，只下载分页内容
        final_md = generate_markdown(
            st.session_state['results'], 
            st.session_state.get('review')
        )
        
        st.download_button(
            label="📥 导出完整讲稿 (.md)",
            data=final_md,
            file_name="presentation_script.md",
            mime="text/markdown"
        )

    # --- 展示区：全局诊断卡片 ---
    if 'review' in st.session_state:
        review = st.session_state['review']
        st.info("💡 **AI 架构师诊断报告**")
        with st.expander("查看详细评价", expanded=True):
            c1, c2 = st.columns(2)
            c1.markdown(f"**🎯 执行摘要**\n\n{review['executive_summary']}")
            c1.markdown(f"**🔗 逻辑流**\n\n{review['logic_diagnosis']}")
            c2.markdown(f"**✅ 亮点**\n\n{review['strengths']}")
            c2.markdown(f"**⚠️ 改进点**\n\n{review['weaknesses']}")
            st.success(f"**🎤 建议结语**: {review['closing_remark']}")

    st.divider()
    
    # --- 展示区：分页内容 ---
    st.caption(f"共分析 {len(st.session_state['results'])} 页")
    for data in st.session_state['results']:
        with st.expander(f"第 {data['index']} 页 | {data.get('visual_summary')}"):
            st.write(data['scripts']['standard'])
